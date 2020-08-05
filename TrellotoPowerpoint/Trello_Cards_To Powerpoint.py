# -*- coding: utf-8 -*-
"""
Created on Thu Jul  2 14:43:17 2020
@author: SPatel9
Ppt of one card 
Have the user pick a board/project to generate a ppt of. 
"""
import requests
from pptx import Presentation 
import os
from pptx.chart.data import CategoryChartData 
from pptx.enum.chart import XL_CHART_TYPE
import tkinter as tk 
from tkinter import *
from tkinter import ttk 
from PIL import ImageTk,Image  
######### Get a Board Title ###############
url = "https://api.trello.com/1/members/me/boards"
headers = {
   "Accept": "application/json"}
#Key and Token for Persoal board 
query = {
   'key': 'b76f5364a94ecd92ef3eae636e59f254', 
   'token': 'd431bd6bbe5191f00b53354f4de4c3a91866c78c94179b0c002a1191502f3b40'}
response = requests.request(
   "GET",
   url,
   headers=headers,
   params=query)

Info = response.json()
### Change bored here ## 
BoardNum=0
board_id=Info[BoardNum]["id"]

######### GET A LIST ############################
url1 = "https://api.trello.com/1/boards/{0}/lists".format(board_id)
response1 = requests.request(
   "GET",
   url1,
   params=query)

ListInfo = response1.json()
list_ids=[]
trello_lists =[]
for y in range(len(ListInfo)):
    list_ids.append(ListInfo[y]["id"]) 
for x in range(len(ListInfo)): 
    trello_lists.append(ListInfo[x]["name"])
    
#################### Get a Card ###############
url_cards=[]
for ids in range(len(list_ids)):
    url_c= "https://api.trello.com/1/lists/{0}/cards".format(list_ids[ids])
    url_cards.append(url_c)

trello_cards=[]
for urlx in url_cards: 
    responses = requests.request(
       "GET",
       urlx,
       params=query
    )
    trello_cards.append(responses.json())

############# Allowing users to pick list and card ###############
colors=['gray','lightsalmon','steelblue']
  
root = tk.Tk()
root.title('Trello Card to PPT')
root.geometry("500x600")
root.wm_iconbitmap("icont2p.ico")

val = StringVar(root)
gal = StringVar(root)

ABOUT_TEXT = """About

TRELLO2PPT will allow users to automatically create powerpoints based on TRELLO Cards.
The more detailed the card is on Trello the better the powerpoint.

There may be some bugs but its my first time making a full program :) 
 
For more information email the creator Soniya Patel 
Soniyaspatel34@gmail.com """

def clickAbout():
    toplevel = Toplevel()
    toplevel.wm_iconbitmap("icont2p.ico")
    label1 = Label(toplevel, text=ABOUT_TEXT, height=0, width=100)
    label1.pack()
numlist = 0
lists=["Select a List"]
while numlist < len(trello_cards):
    lists.append( "{0}".format(trello_lists[numlist]))
    numlist +=1
val.set(lists[0])  
                    
tk.Label(root, 
         text="""Create an Automated Powerpoint based on a Trello Card'""",
         justify = tk.LEFT,
         padx = 20).pack()
#
def comboy2(event):
    Label(root,text="Click 'Create Powerpoint'").pack()
def comboy(event):
    cards= ['Select a Dog']
    cn= 0
    for numcard in trello_cards[(lists.index(val.get()))-1]:
        cards.append("{0}".format(numcard['name']))
        cn+=1 
    gal.set(cards[0])  
    tk.OptionMenu(root,gal,*cards,command=comboy2).pack(side='top',padx=10,pady=10)    
    
tk.OptionMenu(root,val,*lists,command=comboy).pack(side='top',padx=10,pady=10)
image= Image.open("Icon.gif")
photo=ImageTk.PhotoImage(image)
label=Label(root,image=photo)
label.image=photo
label.pack()
create = Button(root, text="Create Powerpoint", width=20, command=root.destroy ,bg = colors[1])
about = Button(root, text="About TRELLO to PPT", width=20, command=clickAbout, bg=colors[2],fg='white')
about.pack(side='bottom',padx=10,pady=10)
create.pack(side='bottom',padx=0,pady=0)
root.mainloop()

l=lists.index(val.get())-1
if l <0:
    exit()
cards= []
cn= 0
for numcard in trello_cards[l]:
    cards.append("{0}".format(numcard['name']))
    cn+=1 
c=cards.index(gal.get())
#### The card picture #####
url_attachments = "https://api.trello.com/1/cards/{0}/attachments".format(trello_cards[l][c]['id'])
responsea = requests.request(
   "GET",
   url_attachments,
   headers=headers,
   params=query
)
attachments = responsea.json()
############Info from cards #############
datelastAct = trello_cards[l][c]['dateLastActivity']
print('Trello Board last updated on: '+datelastAct.split('T')[0])
proj_name= trello_cards[l][c]['name'].replace(':','-')
# Title Slide #
prs = Presentation('0_trello_template.pptx')
#Title Slide 
slide=prs.slides
slide = slide.add_slide( prs.slide_layouts[0])
#slide.shapes.title.text =  'Additive Applications'
slide.placeholders[15].text = 'The Dog House'
slide.placeholders[16].text = proj_name
slide.placeholders[13].text = 'Trello Board last updated on: '+datelastAct.split('T')[0]

# Slide 1: Project Description #
slide1 = prs.slides
slide1 = slide1.add_slide(prs.slide_layouts[1])
slide1.shapes.title.text =  trello_cards[l][c]['name']

####### Table Custom Content #####
#CustomAnswer=[]
#for cusop in range(len(idValue)): 
#    url_cc = "https://api.trello.com/1/customFields/{0}/options".format(idCustomField[cusop])
#    response_ccoptions = requests.request(
#       "GET",
#       url_cc,
#       params=query
#    )
#    CustomAnswer.append(response_ccoptions.json())
#CCIndivAns= []
#ccAnswersiD = []
#j=0
#for lists_cc in CustomAnswer:
#    e=0 
#    ccAnswersiD.append([])
#    while e < len(lists_cc):    
#            ccAnswersiD[j].append(lists_cc[e]['_id'])
#            e+=1
#    j+=1
#i=0
#for sub in ccAnswersiD: 
#    if idValue[i] in sub:
#        CCIndivAns.append(CustomAnswer[ccAnswersiD.index(sub)][sub.index(idValue[i])]['value']['text'])
#        i+=1 
#    
#Fields=[]
#s=0
#for ids in idCustomField: 
#    if idCustomField[s] in CF_dropdownid: 
#        idindex=CF_dropdownid.index(idCustomField[s])
#        Fields.append(Customfields_dropdowns[idindex]['name'])  
#        s+=1

############### Cover Picture ##############
if trello_cards[l][c]['manualCoverAttachment'] == True:
    url_image = attachments[0]['url']
    cover_image = "{0}.jpg".format(proj_name)
    r = requests.get(url_image)
    with open(cover_image,'wb') as file: 
        file.write(r.content)
    slide1.placeholders[1].insert_picture(cover_image)

slide1.placeholders[2].text = trello_cards[l][c]['desc'] # card Description 
slide1.placeholders[15].text = trello_lists[l] # Status

#Division Name 
if len(trello_cards[l][c]['labels']) == 1:
    slide1.placeholders[13].text = trello_cards[l][c]['labels'][0]['name'] # Division 

        
####penning powerpoint #      
proj_name= (trello_cards[l][c]['name'].replace(':','-')).replace('/','&')
prs.save("{0}-Trello.pptx".format(proj_name))
os.startfile("{0}-Trello.pptx".format(proj_name))

